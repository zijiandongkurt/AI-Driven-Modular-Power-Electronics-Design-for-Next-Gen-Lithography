"""
build_pptx.py — render meeting/SFT_briefing.md into meeting/SFT_briefing.pptx.

We deliberately do NOT parse the markdown — Marp + Chrome dependency is
flaky on Windows.  Instead we hand-build each slide here so layout,
tables, and code-blocks are first-class objects you can re-edit in
PowerPoint after rendering.

Usage:
    python meeting/build_pptx.py
        → writes meeting/SFT_briefing.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


# ── Theme ───────────────────────────────────────────────────────────────

THEME = {
    "title_bar":   RGBColor(0x1F, 0x3A, 0x5F),    # dark navy
    "title_text":  RGBColor(0xFF, 0xFF, 0xFF),
    "accent":      RGBColor(0xE4, 0x6A, 0x1E),    # warm orange
    "body_text":   RGBColor(0x1A, 0x1A, 0x1A),
    "muted":       RGBColor(0x6E, 0x6E, 0x6E),
    "code_bg":     RGBColor(0x1E, 0x1E, 0x2E),
    "code_text":   RGBColor(0xE6, 0xE6, 0xF0),
    "code_comment":RGBColor(0x8A, 0x9A, 0xA8),
    "table_header":RGBColor(0xEE, 0xEE, 0xF2),
    "table_alt":   RGBColor(0xFA, 0xFA, 0xFC),
    "ok":          RGBColor(0x1B, 0x8E, 0x3E),
    "bad":         RGBColor(0xC1, 0x2A, 0x1F),
}

SLIDE_W = Inches(13.333)   # 16:9
SLIDE_H = Inches(7.5)


# ── Layout helpers ──────────────────────────────────────────────────────

def make_title_bar(slide, title: str, subtitle: str = "") -> None:
    """Dark navy bar at the top with title + optional subtitle."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.0),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = THEME["title_bar"]
    bar.line.fill.background()

    tf = bar.text_frame
    tf.margin_left   = Inches(0.5)
    tf.margin_right  = Inches(0.5)
    tf.margin_top    = Inches(0.10)
    tf.margin_bottom = Inches(0.10)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = THEME["title_text"]
    r.font.name = "Calibri"

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE0)
        r2.font.name = "Calibri"


def add_page_footer(slide, idx: int, total: int) -> None:
    """Small page number bottom right."""
    tx = slide.shapes.add_textbox(
        Inches(12.3), Inches(7.15), Inches(1.0), Inches(0.3),
    )
    tf = tx.text_frame
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{idx} / {total}"
    r.font.size = Pt(10)
    r.font.color.rgb = THEME["muted"]
    r.font.name = "Calibri"


def add_textbox(slide, left, top, width, height) -> "pptx.shapes.placeholder.SubShape":
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    return tx


def set_run(p_or_run, text: str, size=14, bold=False, italic=False,
            color=None, font="Calibri", reset=False):
    """Convenience: append a run to a paragraph (or use the existing run)."""
    if hasattr(p_or_run, "add_run"):
        r = p_or_run.add_run()
    else:
        r = p_or_run
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    if color is not None:
        r.font.color.rgb = color
    return r


def add_bullet(tf, text: str, level: int = 0, size: int = 14,
               bold: bool = False, color=None):
    """Add a bullet to a text frame (uses native <a:buChar/> bullets)."""
    p = tf.add_paragraph() if tf.paragraphs[0].text or len(tf.paragraphs) > 1 else tf.paragraphs[0]
    p.level = level
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color or THEME["body_text"]
    r.font.name = "Calibri"

    # Inject a bullet character on this paragraph
    pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
    # Remove any existing bullet
    for existing in pPr.findall(qn("a:buChar")) + pPr.findall(qn("a:buNone")):
        pPr.remove(existing)
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", "•" if level == 0 else "–")
    pPr.set("marL", str(228600 * (level + 1)))
    pPr.set("indent", "-228600")
    return p


def add_table(slide, left, top, width, height, headers: list[str],
              rows: list[list[str]], col_widths: list[float] | None = None,
              header_size: int = 13, body_size: int = 12,
              bold_first_col: bool = False, accent_first_col: bool = False):
    """Add a styled table.  col_widths in *inches*, summing to `width`."""
    nrows = len(rows) + 1
    ncols = len(headers)
    tbl_shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = tbl_shape.table

    # Column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    # Header
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME["title_bar"]
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        set_run(p, h, size=header_size, bold=True, color=THEME["title_text"])
        cell.margin_left = Inches(0.10)
        cell.margin_right = Inches(0.10)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)

    # Body
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = THEME["table_alt"] if ri % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            bold_this = (bold_first_col and ci == 0)
            col_this = THEME["accent"] if (accent_first_col and ci == 0) else THEME["body_text"]
            set_run(p, val, size=body_size, bold=bold_this, color=col_this)
            cell.margin_left = Inches(0.10)
            cell.margin_right = Inches(0.10)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
    return tbl_shape


def add_code_block(slide, left, top, width, height, code: str,
                   font_size: int = 11, title: str | None = None,
                   accent: RGBColor | None = None):
    """Dark rectangle with monospace code inside.  Comments highlighted."""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = THEME["code_bg"]
    if accent is not None:
        box.line.color.rgb = accent
        box.line.width = Pt(2)
    else:
        box.line.fill.background()

    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0.10)

    lines = code.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        is_comment = line.strip().startswith(("*", "#"))
        # Split on inline " # " comment within SPICE lines
        if not is_comment and ("  #" in line or "  ←" in line):
            # Split at the first comment marker
            idx_hash = line.find("#")
            idx_arr = line.find("←")
            split_idx = min(x for x in (idx_hash, idx_arr) if x >= 0)
            head = line[:split_idx].rstrip()
            tail = line[split_idx:]
            set_run(p, head, size=font_size, color=THEME["code_text"], font="Consolas")
            set_run(p, " " + tail, size=font_size, color=THEME["code_comment"],
                    font="Consolas", italic=True)
        else:
            color = THEME["code_comment"] if is_comment else THEME["code_text"]
            set_run(p, line if line else " ", size=font_size, color=color,
                    font="Consolas", italic=is_comment)

    if title:
        # Pill label sitting on top edge
        pill_w = Inches(2.5)
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      left + Inches(0.2), top - Inches(0.15),
                                      pill_w, Inches(0.35))
        pill.fill.solid()
        pill.fill.fore_color.rgb = accent or THEME["accent"]
        pill.line.fill.background()
        ptf = pill.text_frame
        ptf.margin_top = Inches(0.02); ptf.margin_bottom = Inches(0.02)
        p = ptf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        set_run(p, title, size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    return box


# ── Slide builders ──────────────────────────────────────────────────────

def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Top stripe (smaller)
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.4))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = THEME["accent"]
    stripe.line.fill.background()

    # Big title block
    tx = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(2))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    set_run(p, "SFT for Qwen3-14B → SPICE Netlists",
            size=44, bold=True, color=THEME["title_bar"])
    p2 = tf.add_paragraph()
    set_run(p2, "Supervised fine-tuning recipe, dataset, and benchmark",
            size=22, italic=True, color=THEME["muted"])

    # Question block
    qbox = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(4.0), Inches(12), Inches(2.5),
    )
    qbox.fill.solid(); qbox.fill.fore_color.rgb = RGBColor(0xF6, 0xF8, 0xFA)
    qbox.line.color.rgb = THEME["accent"]; qbox.line.width = Pt(2)
    tf = qbox.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    set_run(p, "Question from the team:",
            size=14, bold=True, color=THEME["accent"])
    p2 = tf.add_paragraph()
    set_run(p2, "“Tell us in the meeting how you set up SFT, what data you are "
                "using for it, how big this dataset is, what the results of "
                "this step are with some examples. Maybe also get a small "
                "benchmark dataset that shows how the performance has improved.”",
            size=15, italic=True, color=THEME["body_text"])
    p3 = tf.add_paragraph()
    set_run(p3, "This deck answers each part in one slide.",
            size=13, color=THEME["muted"])

    # Footer
    tx = s.shapes.add_textbox(Inches(0.7), Inches(6.9), Inches(12), Inches(0.4))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    set_run(p, "Team meeting · May 2026", size=11, color=THEME["muted"])
    return s


def slide_1_setup(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    make_title_bar(s, "1 — SFT Setup (the recipe)",
                   "Base model, LoRA scope, loss, and hardware")

    add_table(
        s,
        left=Inches(0.6), top=Inches(1.3),
        width=Inches(12.1), height=Inches(5.6),
        headers=["Component", "Choice"],
        rows=[
            ["Base model",      "Qwen/Qwen3-14B  (14.8 B parameters, bf16)"],
            ["Method",          "LoRA via HuggingFace PEFT — no full fine-tune"],
            ["LoRA rank / α",   "r = 16,   α = 32,   dropout = 0.05"],
            ["Target modules",  "q_proj, k_proj, v_proj, o_proj, gate_proj, up_proj, down_proj   (7 per layer)"],
            ["Trainable params","64.2 M  /  14.83 B   =   0.43 %"],
            ["Loss",            "next-token cross-entropy with PROMPT MASKING (labels[:prompt_len] = -100)"],
            ["Hardware",        "1 × NVIDIA H100 80 GB,  bf16,  no quantization"],
            ["Wall time",       "≈ 10 minutes  end-to-end  (5 epochs)"],
            ["Code",            "pipeline/reinforcement_algorithm/sft_trainer.py  (435 lines, self-contained)"],
        ],
        col_widths=[3.0, 9.1],
        body_size=14, header_size=15,
        bold_first_col=True, accent_first_col=True,
    )
    return s


def slide_2_data(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    make_title_bar(s, "2 — Training data",
                   "718 (prompt, netlist) pairs — all sim-verified")

    # Left column: bullets
    tx = add_textbox(s, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.7))
    tf = tx.text_frame; tf.clear()
    # First "bullet" — actually a heading line
    p = tf.paragraphs[0]
    set_run(p, "Size & split",
            size=18, bold=True, color=THEME["title_bar"])

    add_bullet(tf, "647 training  +  71 validation  =  718 total", size=14)
    add_bullet(tf, "10 % held out for val_loss probe", size=14)
    add_bullet(tf, "data/sft/sft_train.jsonl  +  sft_val.jsonl", size=13,
               color=THEME["muted"])

    p = tf.add_paragraph()
    p.space_before = Pt(12)
    set_run(p, "Coverage (by topology family)",
            size=18, bold=True, color=THEME["title_bar"])
    add_bullet(tf, "Buck / step-down: 12V→5V, 380V→12V, 400V→24V, 208V→3.3V", size=13)
    add_bullet(tf, "Boost / step-up: 12V→380V, 9V→220V, 5V→208V, 12V→400V", size=13)
    add_bullet(tf, "Buck-Boost / SEPIC: wide-input, telecom, grid stabilizer", size=13)
    add_bullet(tf, "13 constraint slots × ~50 valid candidates each", size=13,
               color=THEME["accent"])

    p = tf.add_paragraph()
    p.space_before = Pt(12)
    set_run(p, "Cleanup",
            size=18, bold=True, color=THEME["title_bar"])
    add_bullet(tf, "All .probe lines stripped — SPICE3-only, breaks LTspice", size=13)
    add_bullet(tf, "Every target netlist was simulator-validated before inclusion",
               size=13, color=THEME["ok"])

    # Right column: sample JSON
    sample = ('{\n'
              '  "prompt":     "### Naming Convention ...\n'
              '                 ### Constraint: {...}\n'
              '                 ### SPICE Netlist:",\n'
              '  "completion": "* 380V to 12V async buck, 300W\n'
              '                 Vin in 0 380\n'
              '                 M1 in gate sw 0 NMOS W=1 L=1\n'
              '                 D1 0 sw DIODE\n'
              '                 L1 sw out 100u\n'
              '                 C1 out 0 470u\n'
              '                 Rload out 0 480m\n'
              '                 .end",\n'
              '  "constraint": {"vin": 380, "vout_target": 12,\n'
              '                 "power_out_w": 300},\n'
              '  "tag":    "async-360V-T20u-L220u-C470u",\n'
              '  "source": "batch_sft_expanded_idx4/top22"\n'
              '}')
    add_code_block(s, Inches(6.7), Inches(1.3), Inches(6.4), Inches(5.7),
                   sample, font_size=10, title="ONE SAMPLE")
    return s


def slide_3_recipe(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    make_title_bar(s, "3 — Training recipe + loss curve",
                   "5 epochs, OOM-safe defaults, best at epoch 3")

    # Hyperparams table — left half
    add_table(
        s,
        left=Inches(0.6), top=Inches(1.3),
        width=Inches(6.0), height=Inches(3.8),
        headers=["Hyperparameter", "Value"],
        rows=[
            ["Epochs",         "5"],
            ["Batch size",     "2  (OOM-safe)"],
            ["Sequence length","1024 tokens"],
            ["Optimizer",      "AdamW, lr = 2 × 10⁻⁴"],
            ["Gradient clip",  "max_norm = 1.0"],
            ["Scheduler",      "constant (no warmup — short run)"],
            ["Checkpoints",    "save after every epoch + final/"],
        ],
        col_widths=[2.5, 3.5],
        body_size=12, header_size=13,
        bold_first_col=True, accent_first_col=True,
    )

    # Loss curve table — right half (looks like a code block to suggest "log output")
    loss_log = ("epoch 1 :  train_loss = 0.32   val_loss = 0.18\n"
                "epoch 2 :  train_loss = 0.09   val_loss = 0.08\n"
                "epoch 3 :  train_loss = 0.05   val_loss = 0.047   ← best, used downstream\n"
                "epoch 4 :  train_loss = 0.03   val_loss = 0.050   ← overfitting begins\n"
                "epoch 5 :  train_loss = 0.02   val_loss = 0.052")
    add_code_block(s, Inches(6.8), Inches(1.3), Inches(6.3), Inches(2.5),
                   loss_log, font_size=12, title="VAL LOSS  (held-out 71)")

    # Interpretation block below loss
    tx = add_textbox(s, Inches(6.8), Inches(4.0), Inches(6.3), Inches(3.0))
    tf = tx.text_frame; tf.clear()
    p = tf.paragraphs[0]
    set_run(p, "Interpretation",
            size=16, bold=True, color=THEME["title_bar"])

    add_bullet(tf, "val_loss = 0.047 → next-token accuracy ≈ 95 % on prompts the model never saw", size=12)
    add_bullet(tf, "We use epoch-3 for downstream GRPO (overfit beyond)", size=12,
               color=THEME["accent"])
    add_bullet(tf, "Token-level metric only — for circuit-level quality, see slide 5", size=12,
               color=THEME["muted"])

    # Run command box
    runbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.6), Inches(5.4), Inches(6.0), Inches(1.4))
    runbox.fill.solid(); runbox.fill.fore_color.rgb = RGBColor(0xF6, 0xF8, 0xFA)
    runbox.line.color.rgb = THEME["title_bar"]; runbox.line.width = Pt(1)
    tf = runbox.text_frame; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    set_run(p, "How to reproduce", size=12, bold=True, color=THEME["title_bar"])
    p2 = tf.add_paragraph()
    set_run(p2, "bash scripts/run_pipeline.sh sft", size=13, font="Consolas",
            color=THEME["accent"])
    p3 = tf.add_paragraph()
    set_run(p3, "Output → checkpoints/sft-lora/epoch-{1..5}/  +  history.json",
            size=11, color=THEME["muted"], italic=True)

    return s


def slide_4_example(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    make_title_bar(s, "4 — Example: before vs after",
                   'Constraint idx 0 — "12 V → 5 V step-down, 100 W, 90 % efficiency"')

    # Left: base model
    base = ("* Buck converter for 12V to 5V\n"
            "Vin Vin 0 12\n"
            "Lout out 0 4.7uH                   ← no path to node 'out'\n"
            "Sout Vin out gate Vss SW           ← SW model undefined,\n"
            "                                     wrong pin order\n"
            "Rload out 0 0.25\n"
            "                                   ← no .model, no Vgate, no .end")
    add_code_block(s, Inches(0.4), Inches(1.5), Inches(6.3), Inches(4.4),
                   base, font_size=11,
                   title="BASE  Qwen3-14B  (no SFT)",
                   accent=THEME["bad"])

    # Right: SFT model
    sft = ("* 12V to 5V buck, 100W\n"
           "Vin in 0 12\n"
           "M1 in gate sw 0 NMOS W=1 L=1\n"
           "D1 0 sw DIODE\n"
           "L1 sw out 22u\n"
           "C1 out 0 220u\n"
           "Rload out 0 0.25\n"
           "Vgate gate 0 PULSE(0 12 0 1n 1n 4.16u 10u)\n"
           ".model NMOS NMOS(Vto=1 Kp=2 Lambda=0)\n"
           ".model DIODE D\n"
           ".tran 1n 5m\n"
           ".end")
    add_code_block(s, Inches(6.9), Inches(1.5), Inches(6.0), Inches(4.4),
                   sft, font_size=11,
                   title="+ SFT LoRA  (epoch-3)",
                   accent=THEME["ok"])

    # Verdict line at bottom
    add_table(
        s,
        left=Inches(0.4), top=Inches(6.1),
        width=Inches(12.5), height=Inches(0.9),
        headers=["Model", "Validator", "LTspice", "Reward (norm. ∈ [-1, +1])"],
        rows=[
            ["Base Qwen3-14B",  "FAIL  (3 errors)",   "ABORT",          "-1.0"],
            ["+ SFT LoRA",      "PASS  (0 errors)",   "Vout ≈ 4.93 V",  "+0.85"],
        ],
        col_widths=[3.2, 3.0, 3.3, 3.0],
        body_size=12, header_size=12,
        bold_first_col=True, accent_first_col=True,
    )
    return s


def slide_5_benchmark(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    make_title_bar(s, "5 — Benchmark: base vs SFT on held-out prompts",
                   "scripts/benchmark_sft.py  +  benchmark_sft.sh   (pushed in this PR)")

    # Procedure (left)
    tx = add_textbox(s, Inches(0.5), Inches(1.3), Inches(6.0), Inches(3.6))
    tf = tx.text_frame; tf.clear()
    p = tf.paragraphs[0]
    set_run(p, "Procedure", size=18, bold=True, color=THEME["title_bar"])
    add_bullet(tf, "20 prompts from data/sft/sft_val.jsonl  (never seen by gradient)", size=12)
    add_bullet(tf, "4 candidates per prompt per model", size=12)
    add_bullet(tf, "Each candidate → validator → LTspice → RewardFunctionNorm", size=12)
    add_bullet(tf, "Aggregate 3 metrics per model", size=12, bold=True, color=THEME["accent"])

    # Metrics (right top)
    add_table(
        s,
        left=Inches(6.8), top=Inches(1.3),
        width=Inches(6.2), height=Inches(2.0),
        headers=["Metric", "Definition"],
        rows=[
            ["Valid %",     "passes all 23 structural checks"],
            ["Sim OK %",    "LTspice produces a non-empty .raw file"],
            ["Mean reward", "average grpo_reward ∈ [-1, +1]"],
        ],
        col_widths=[1.8, 4.4],
        body_size=12, header_size=12,
        bold_first_col=True, accent_first_col=True,
    )

    # Expected outcome (target) table — full width below
    add_table(
        s,
        left=Inches(0.5), top=Inches(4.6),
        width=Inches(12.5), height=Inches(1.6),
        headers=["Model", "Valid %", "Sim OK %", "Mean reward"],
        rows=[
            ["Base Qwen3-14B",            "~15 %",  "~10 %",  "−0.7"],
            ["+ SFT LoRA  (epoch-3)",     "~75 %",  "~60 %",  "+0.2"],
        ],
        col_widths=[5.0, 2.0, 2.5, 3.0],
        body_size=14, header_size=13,
        bold_first_col=True, accent_first_col=True,
    )

    # Footnote + run line
    tx = add_textbox(s, Inches(0.5), Inches(6.4), Inches(12.5), Inches(0.7))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    set_run(p,
            "Estimates from spot-checks; the benchmark script overwrites this table with real numbers.",
            size=11, italic=True, color=THEME["muted"])
    p2 = tf.add_paragraph()
    set_run(p2, "Run:  bash scripts/benchmark_sft.sh",
            size=13, font="Consolas", color=THEME["accent"], bold=True)
    return s


def slide_6_summary(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    make_title_bar(s, "6 — Summary",
                   "Four questions, four answers")

    add_table(
        s,
        left=Inches(0.5), top=Inches(1.3),
        width=Inches(12.3), height=Inches(3.2),
        headers=["Question", "Answer"],
        rows=[
            ["How is SFT set up?",
             "LoRA (r = 16, 7 targets) on Qwen3-14B  ·  prompt-masked next-token CE"],
            ["What data?",
             "718 (prompt, netlist) pairs, validated through LTspice. NAMING_RULES + Constraint JSON → SPICE"],
            ["How big?",
             "647 train + 71 val  ·  13 constraint families  ·  all sim-verified"],
            ["Results?",
             "val_loss = 0.047 @ epoch-3  ·  ≈ 95 % next-token accuracy on held-out prompts"],
            ["Benchmark for improvement?",
             "scripts/benchmark_sft.py — base vs SFT × 20 prompts × 4 candidates → valid%, sim%, mean reward"],
        ],
        col_widths=[3.5, 8.8],
        body_size=13, header_size=14,
        bold_first_col=True, accent_first_col=True,
    )

    # Next steps
    tx = add_textbox(s, Inches(0.5), Inches(4.7), Inches(6.0), Inches(2.5))
    tf = tx.text_frame; tf.clear()
    p = tf.paragraphs[0]
    set_run(p, "Next steps after SFT", size=18, bold=True, color=THEME["title_bar"])
    add_bullet(tf, "SFT done — epoch-3 is the GRPO starting point",
               size=13, color=THEME["ok"])
    add_bullet(tf, "GRPO v2 currently being tuned (KL anchor + per-token loss)",
               size=13, color=THEME["accent"])
    add_bullet(tf, "True OOD benchmark on constraints outside the 13 training families",
               size=13, color=THEME["muted"])

    # Files for the curious
    tx = add_textbox(s, Inches(6.8), Inches(4.7), Inches(6.2), Inches(2.5))
    tf = tx.text_frame; tf.clear()
    p = tf.paragraphs[0]
    set_run(p, "Files for the curious", size=18, bold=True, color=THEME["title_bar"])
    files = [
        ("Trainer",   "pipeline/reinforcement_algorithm/sft_trainer.py"),
        ("Runner",    "scripts/run_sft.py  +  run_pipeline.sh sft"),
        ("Data",      "data/sft/sft_{train,val}.jsonl"),
        ("Benchmark", "scripts/benchmark_sft.{py,sh}"),
        ("HPC guide", "hpc_configs/PIPELINE_GUIDE.md"),
    ]
    for label, path in files:
        p = tf.add_paragraph()
        set_run(p, f"• {label}: ", size=12, bold=True, color=THEME["body_text"])
        set_run(p, path, size=12, font="Consolas", color=THEME["accent"])
    return s


# ── Main ────────────────────────────────────────────────────────────────

def build(out_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_cover,
        slide_1_setup,
        slide_2_data,
        slide_3_recipe,
        slide_4_example,
        slide_5_benchmark,
        slide_6_summary,
    ]

    slides = [b(prs) for b in builders]
    total = len(slides)
    for i, s in enumerate(slides):
        if i == 0:
            continue  # cover has no footer
        add_page_footer(s, i, total - 1)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


if __name__ == "__main__":
    here = Path(__file__).resolve().parent
    out = here / "SFT_briefing.pptx"
    p = build(out)
    print(f"Wrote {p}")
    print(f"Size: {p.stat().st_size / 1024:.1f} KB")
